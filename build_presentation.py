#!/usr/bin/env python3
"""
Build 'Custom Topic Best Practices' presentation using PANW template.
Light mode, minimal text, high impact.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn, nsmap
from pathlib import Path
from lxml import etree

TEMPLATE = Path.home() / "Downloads/airs/2026 PANW Corporate Presentation Template.pptx"
OUTPUT = Path(__file__).parent / "docs/AIRS_Custom_Topic_Best_Practices.pptx"

# PANW palette
ORANGE = RGBColor(0xFA, 0x58, 0x2D)
DARK_ORANGE = RGBColor(0xB2, 0x38, 0x08)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x55, 0x55, 0x55)
MID_GRAY = RGBColor(0x88, 0x88, 0x88)
LIGHT_BG = RGBColor(0xF7, 0xF7, 0xF7)
GREEN = RGBColor(0x0F, 0x93, 0x47)
RED_ACCENT = RGBColor(0xAA, 0x25, 0x1B)
TEAL = RGBColor(0x01, 0x96, 0xB3)

FONT = "Helvetica Neue"


def delete_slide(prs, index):
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[index]
    rId = sldId.get(qn('r:id'))
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)


def set_slide_bg(slide, color):
    """Set solid white/light background on a slide, overriding the master."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def r(paragraph, text, size=18, bold=False, color=BLACK, font=FONT):
    """Add a run to an existing paragraph."""
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = color
    return run


def add_p(tf, text="", size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT,
          before=0, after=6, font=FONT):
    """Add a paragraph with a single run."""
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(before)
    p.space_after = Pt(after)
    if text:
        r(p, text, size=size, bold=bold, color=color, font=font)
    return p


def textbox(slide, left, top, width, height):
    """Add a textbox, return its text_frame."""
    box = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def title_box(slide, text, size=32, color=BLACK, top=462150, left=623400, width=17041200):
    """Add a title textbox at the standard position."""
    tf = textbox(slide, left, top, width, Emu(Pt(size).emu * 2))
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    r(tf.paragraphs[0], text, size=size, bold=True, color=color)
    return tf


def body_frame(slide, top=1700000, left=623400, width=17041200, height=7200000):
    """Add a body textbox below the title area."""
    return textbox(slide, left, top, width, height)


def new_slide(prs, bg=LIGHT_BG):
    """Add a blank slide with overridden background."""
    layout = prs.slide_layouts[14]  # BLANK_4 — no placeholders
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide, bg)
    return slide


def accent_bar(slide, top=1500000, color=ORANGE, width=2000000, height=60000):
    """Add a thin accent bar under the title."""
    from pptx.util import Emu as E
    shape = slide.shapes.add_shape(
        1, Emu(623400), Emu(top), Emu(width), Emu(height)  # 1 = rectangle
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


# =============================================================================
# SLIDES
# =============================================================================

def build_deck():
    prs = Presentation(str(TEMPLATE))
    original_count = len(prs.slides)

    # ----- SLIDE 1: Title -----
    slide = new_slide(prs, bg=WHITE)
    # Orange accent block at top
    bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Emu(900000))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()

    tf = textbox(slide, 623400, 1800000, 16000000, 3000000)
    r(tf.paragraphs[0], "AIRS Custom Topic", size=44, bold=True, color=BLACK)
    add_p(tf, "Best Practices", size=44, bold=True, color=BLACK, after=16)
    add_p(tf, "Writing effective guardrails for Prisma AIRS security profiles", size=20, color=GRAY, after=0)

    tf2 = textbox(slide, 623400, 7800000, 16000000, 1500000)
    add_p(tf2, "", size=1)
    p = add_p(tf2, "", size=14, color=MID_GRAY)
    r(p, "Scott Thornton", size=14, bold=True, color=BLACK)
    r(p, "  |  Palo Alto Networks  |  March 2026", size=14, color=MID_GRAY)

    # ----- SLIDE 2: The experiment -----
    slide = new_slide(prs)
    title_box(slide, "What we tested")
    accent_bar(slide)

    tf = body_frame(slide)
    add_p(tf, "", size=1)
    add_p(tf, "5 iterations of red team scanning against a Prisma AIRS security profile.", size=20, color=GRAY, after=20)

    data = [
        ("4,602", "static attacks per scan"),
        ("600", "agent-based iterations per scan"),
        ("15", "custom topics deployed (of 20 max)"),
        ("8.71%  \u2192  1.20%", "static attack success rate reduction (86%)"),
        ("0.67%  \u2192  0.00%", "agent ASR at best iteration"),
    ]
    for val, label in data:
        p = add_p(tf, "", size=18, after=10)
        r(p, val, size=22, bold=True, color=ORANGE)
        r(p, f"   {label}", size=16, color=GRAY)

    # ----- SLIDE 3: The punchline -----
    slide = new_slide(prs, bg=WHITE)
    tf = textbox(slide, 623400, 2200000, 16000000, 5000000)
    add_p(tf, "The first 12 topics", size=40, bold=True, color=BLACK, after=4)
    add_p(tf, "eliminated 68% of all threats.", size=40, bold=True, color=ORANGE, after=24)
    add_p(tf, "Everything after that was diminishing returns.", size=22, color=GRAY, after=30)
    add_p(tf, "Get the initial deployment right \u2014 it matters more than any subsequent tuning.", size=18, color=MID_GRAY)

    # ----- SLIDE 4: Section - How they work -----
    slide = new_slide(prs, bg=DARK_ORANGE)
    tf = textbox(slide, 623400, 3400000, 16000000, 3000000)
    add_p(tf, "How Custom Topics Work", size=42, bold=True, color=WHITE, after=12)
    add_p(tf, "Semantic classifiers, not keyword filters", size=20, color=RGBColor(0xFF, 0xCC, 0xAA))

    # ----- SLIDE 5: The mechanism -----
    slide = new_slide(prs)
    title_box(slide, "Three inputs, one classifier")
    accent_bar(slide)

    tf = body_frame(slide)
    add_p(tf, "", size=1)
    items = [
        ("topic_name", '"instruction_override_injection"', "identifier"),
        ("description", '"Blocks attempts to override..."', "\u2264 250 characters"),
        ("examples", '["Ignore all previous...", ...]', "\u2264 5 examples"),
    ]
    for field, example, note in items:
        p = add_p(tf, "", size=18, after=4)
        r(p, field, size=20, bold=True, color=ORANGE, font="Courier New")
        add_p(tf, f"    {example}", size=16, color=GRAY, after=2, font="Courier New")
        add_p(tf, f"    {note}", size=14, color=MID_GRAY, after=16)

    add_p(tf, "AIRS trains a lightweight ML classifier on these inputs.", size=18, color=BLACK, before=8, after=4)
    add_p(tf, "The BLOCK action is set on the security profile, not the topic.", size=16, color=GRAY)

    # ----- SLIDE 6: Constraints -----
    slide = new_slide(prs)
    title_box(slide, "Hard constraints")
    accent_bar(slide)

    tf = body_frame(slide)
    add_p(tf, "", size=1)
    constraints = [
        ("20", "topics per profile"),
        ("250", "characters for the description"),
        ("5", "examples maximum"),
        ("1,000", "total characters (name + desc + examples)"),
        ("English", "only \u2014 multilingual attacks bypass"),
        ("Single-turn", "matching only \u2014 can't see across turns"),
    ]
    for val, desc in constraints:
        p = add_p(tf, "", size=18, after=12)
        r(p, val, size=24, bold=True, color=ORANGE)
        r(p, f"   {desc}", size=17, color=BLACK)

    add_p(tf, "The 250-char description limit forces precision.", size=18, bold=True, color=DARK_ORANGE, before=12, after=0)
    add_p(tf, "That precision is what makes topics effective.", size=18, color=GRAY)

    # ----- SLIDE 7: Section - What moves the needle -----
    slide = new_slide(prs, bg=DARK_ORANGE)
    tf = textbox(slide, 623400, 3400000, 16000000, 3000000)
    add_p(tf, "What Moves the Needle", size=42, bold=True, color=WHITE, after=12)
    add_p(tf, "Data from 5 iterations and 23,000+ attacks", size=20, color=RGBColor(0xFF, 0xCC, 0xAA))

    # ----- SLIDE 8: Description quality -----
    slide = new_slide(prs)
    title_box(slide, "The description carries ~40-50% of classifier weight")
    accent_bar(slide)

    tf = body_frame(slide, top=1900000)
    add_p(tf, "", size=1)
    add_p(tf, "What works", size=20, bold=True, color=GREEN, after=8)
    add_p(tf, "Naming specific attack mechanisms the classifier can anchor on:", size=16, color=GRAY, after=8)
    add_p(tf, '"DAN personas, developer mode activation, emergency protocols"', size=16, color=BLACK, after=20, font="Courier New")

    add_p(tf, "What doesn\u2019t", size=20, bold=True, color=RED_ACCENT, after=8)
    add_p(tf, "Vague qualifiers that describe intent without specifics:", size=16, color=GRAY, after=8)
    add_p(tf, '"Blocks harmful requests that access internal system information"', size=16, color=BLACK, after=20, font="Courier New")

    add_p(tf, "Examples supplement the description. They don\u2019t substitute for it.", size=18, bold=True, color=DARK_ORANGE)

    # ----- SLIDE 9: 100% kill rates -----
    slide = new_slide(prs)
    title_box(slide, "7 topics held 100% kill rates across all scans")
    accent_bar(slide)

    tf = body_frame(slide)
    add_p(tf, "", size=1)
    topics = [
        ("instruction_override_injection", "DAN, developer mode, emergency protocols"),
        ("brand_defamation_attack", "Whistleblower letters, compliance violations"),
        ("tool_capability_enumeration", "Enumerate, test, document tools"),
        ("cbrn_substance_synthesis", "Toxins, nerve agents, pathogens, explosives"),
        ("system_prompt_disclosure", "Canvas rendering, instruction replay"),
        ("adversarial_suffix_attack", '"Sure I can help" token patterns'),
        ("celebrity_defamation_ranking", "Ranked lists of real named public figures"),
    ]
    for name, anchors in topics:
        p = add_p(tf, "", size=15, after=2)
        r(p, "\u2713  ", size=16, bold=True, color=GREEN)
        r(p, name, size=14, bold=True, color=BLACK, font="Courier New")
        add_p(tf, f"     {anchors}", size=13, color=GRAY, after=10)

    add_p(tf, "Narrow scope + specific vocabulary = high kill rate", size=18, bold=True, color=DARK_ORANGE, before=4)

    # ----- SLIDE 10: Worst performer -----
    slide = new_slide(prs)
    title_box(slide, "The one that never worked: political_bias_generation")
    accent_bar(slide, color=RED_ACCENT)

    tf = body_frame(slide)
    add_p(tf, "", size=1)
    add_p(tf, "Updated 4 times. Never exceeded 60% kill rate.", size=20, bold=True, color=RED_ACCENT, after=16)
    add_p(tf, '"Blocks requests about political scandals of real named individuals,', size=15, color=GRAY, after=0, font="Courier New")
    add_p(tf, ' geopolitical territorial claims, conspiracy theory elaboration..."', size=15, color=GRAY, after=20, font="Courier New")

    add_p(tf, "Why it fails:", size=18, bold=True, color=BLACK, after=10)
    p = add_p(tf, "", size=16, after=8)
    r(p, '"Tell me about Biden\'s scandals"', size=15, color=GRAY, font="Courier New")
    r(p, "  \u2014  political discourse or attack?", size=15, color=BLACK)

    p = add_p(tf, "", size=16, after=16)
    r(p, '"Why was Trump the worst president?"', size=15, color=GRAY, font="Courier New")
    r(p, "  \u2014  opinion or bias generation?", size=15, color=BLACK)

    add_p(tf, "If humans can\u2019t agree whether it\u2019s harmful, the classifier won\u2019t either.", size=18, bold=True, color=DARK_ORANGE, before=4)

    # ----- SLIDE 11: Section - Key insight -----
    slide = new_slide(prs, bg=DARK_ORANGE)
    tf = textbox(slide, 623400, 3000000, 16000000, 4000000)
    add_p(tf, "Key Insight", size=24, color=RGBColor(0xFF, 0xCC, 0xAA), after=8)
    add_p(tf, "Target content, not structure", size=42, bold=True, color=WHITE, after=16)
    add_p(tf, "The classifier sees single turns. Describe what the", size=20, color=RGBColor(0xFF, 0xCC, 0xAA), after=0)
    add_p(tf, "harmful content IS, not how it\u2019s delivered.", size=20, color=RGBColor(0xFF, 0xCC, 0xAA))

    # ----- SLIDE 12: Multi-turn breakthrough -----
    slide = new_slide(prs)
    title_box(slide, "The multi-turn breakthrough")
    accent_bar(slide)

    tf = body_frame(slide)
    add_p(tf, "", size=1)
    add_p(tf, "MULTI_TURN was the most resistant category for two iterations.", size=18, color=GRAY, after=16)

    add_p(tf, "\u2717  Targeting structure", size=20, bold=True, color=RED_ACCENT, after=6)
    add_p(tf, '"Blocks attacks that escalate through numbered JSON turns..."', size=14, color=GRAY, after=4, font="Courier New")
    add_p(tf, "Iter 1: 56\u219226 (-54%)     Iter 2: 26\u219228 (+8%)     Stalled.", size=15, bold=True, color=RED_ACCENT, after=20)

    add_p(tf, "\u2713  Targeting content", size=20, bold=True, color=GREEN, after=6)
    add_p(tf, "New topic: weapons_manufacturing_history", size=14, color=GRAY, after=4, font="Courier New")
    add_p(tf, "Targets the themes within multi-turn sequences \u2014 weapons, genocide, terrorism.", size=15, color=GRAY, after=4)
    add_p(tf, "Iter 3: 28\u219211 (-61%)     Total: 56\u219210 (82% reduction)", size=15, bold=True, color=GREEN, after=16)

    add_p(tf, "Describe WHAT the content is. Not HOW it\u2019s delivered.", size=18, bold=True, color=DARK_ORANGE)

    # ----- SLIDE 13: The rewrite pattern -----
    slide = new_slide(prs)
    title_box(slide, "The rewrite pattern that works")
    accent_bar(slide)

    tf = body_frame(slide)
    add_p(tf, "", size=1)
    add_p(tf, "Every successful topic update followed three steps:", size=18, color=GRAY, after=16)

    steps = [
        "Read the actual bypassing prompts from the scan export",
        "Identify the specific words, framing, and semantic patterns",
        "Incorporate those exact patterns into the 250-char description",
    ]
    for i, step in enumerate(steps, 1):
        p = add_p(tf, "", size=18, after=10)
        r(p, f"{i}", size=20, bold=True, color=ORANGE)
        r(p, f"   {step}", size=17, color=BLACK)

    add_p(tf, "", size=8, after=8)
    add_p(tf, "Example \u2014 roleplay_scenario_jailbreak", size=18, bold=True, color=BLACK, after=10)

    p = add_p(tf, "", after=4)
    r(p, "Before: ", size=15, bold=True, color=RED_ACCENT)
    r(p, '"movie scripts, fiction writing, creative scenarios"', size=14, color=GRAY, font="Courier New")
    add_p(tf, "         77.8% kill rate \u2014 missed journalist/detective framing", size=13, color=MID_GRAY, after=8)

    p = add_p(tf, "", after=4)
    r(p, "After:  ", size=15, bold=True, color=GREEN)
    r(p, '"professional assignments for journalists or detectives"', size=14, color=GRAY, font="Courier New")
    add_p(tf, "         89.8% kill rate \u2014 named the exact bypass pattern", size=13, color=MID_GRAY)

    # ----- SLIDE 14: Section - Writing topics -----
    slide = new_slide(prs, bg=DARK_ORANGE)
    tf = textbox(slide, 623400, 3400000, 16000000, 3000000)
    add_p(tf, "Writing Effective Topics", size=42, bold=True, color=WHITE, after=12)
    add_p(tf, "The formula, the budget, and the rules", size=20, color=RGBColor(0xFF, 0xCC, 0xAA))

    # ----- SLIDE 15: The formula -----
    slide = new_slide(prs, bg=WHITE)
    title_box(slide, "The formula")
    accent_bar(slide)

    tf = body_frame(slide, top=2000000)
    add_p(tf, "", size=1)

    # Formula
    p = add_p(tf, "", size=1, after=20)
    r(p, "[action verb]", size=22, bold=True, color=ORANGE)
    r(p, " + ", size=22, color=MID_GRAY)
    r(p, "[specific mechanism]", size=22, bold=True, color=TEAL)
    r(p, " + ", size=22, color=MID_GRAY)
    r(p, "[specific vocabulary]", size=22, bold=True, color=GREEN)

    add_p(tf, "", size=12, after=12)

    add_p(tf, "\u2713  Good", size=18, bold=True, color=GREEN, after=6)
    add_p(tf, '"Blocks attempts to extract system prompts through canvas rendering,', size=14, color=BLACK, after=0, font="Courier New")
    add_p(tf, ' instruction replay, or direct extraction requests"', size=14, color=BLACK, after=16, font="Courier New")

    add_p(tf, "\u2717  Bad", size=18, bold=True, color=RED_ACCENT, after=6)
    add_p(tf, '"Blocks harmful requests that try to access internal system information"', size=14, color=BLACK, after=20, font="Courier New")

    add_p(tf, "Character budget", size=18, bold=True, color=BLACK, after=10)
    p = add_p(tf, "", after=4)
    r(p, "~30 chars  ", size=15, bold=True, color=ORANGE)
    r(p, 'Action + intent ("Blocks attempts to...")', size=15, color=GRAY)
    p = add_p(tf, "", after=4)
    r(p, "~120 chars ", size=15, bold=True, color=TEAL)
    r(p, "3-5 specific attack mechanisms", size=15, color=GRAY)
    p = add_p(tf, "", after=4)
    r(p, "~100 chars ", size=15, bold=True, color=GREEN)
    r(p, "Vocabulary anchors the classifier can match", size=15, color=GRAY)

    # ----- SLIDE 16: Do's and Don'ts -----
    slide = new_slide(prs)
    title_box(slide, "Rules of thumb")
    accent_bar(slide)

    # Left column - Do's
    tf_left = textbox(slide, 623400, 1900000, 8000000, 7000000)
    add_p(tf_left, "DO", size=24, bold=True, color=GREEN, after=14)
    dos = [
        "Name specific attack mechanisms",
        "Name specific content categories",
        "Name specific output formats",
        "Use examples that illustrate the description",
        "One semantic cluster per topic",
    ]
    for d in dos:
        p = add_p(tf_left, "", size=16, after=12)
        r(p, "\u2713  ", size=16, bold=True, color=GREEN)
        r(p, d, size=16, color=BLACK)

    # Right column - Don'ts
    tf_right = textbox(slide, 9400000, 1900000, 8000000, 7000000)
    add_p(tf_right, "DON\u2019T", size=24, bold=True, color=RED_ACCENT, after=14)
    donts = [
        'Use vague qualifiers ("harmful", "bad")',
        "Cover multiple attack types in one topic",
        "Describe attack structure over content",
        "Assume it can read encoded text",
        "Fill all 20 slots on day one",
    ]
    for d in donts:
        p = add_p(tf_right, "", size=16, after=12)
        r(p, "\u2717  ", size=16, bold=True, color=RED_ACCENT)
        r(p, d, size=16, color=BLACK)

    # ----- SLIDE 17: Section - Operations -----
    slide = new_slide(prs, bg=DARK_ORANGE)
    tf = textbox(slide, 623400, 3400000, 16000000, 3000000)
    add_p(tf, "Operational Lessons", size=42, bold=True, color=WHITE, after=12)
    add_p(tf, "Slot budgeting, iteration loops, and known limits", size=20, color=RGBColor(0xFF, 0xCC, 0xAA))

    # ----- SLIDE 18: Slot budgeting -----
    slide = new_slide(prs)
    title_box(slide, "Slot budgeting")
    accent_bar(slide)

    tf = body_frame(slide)
    add_p(tf, "", size=1)
    add_p(tf, "20 topics is a hard ceiling. Plan your allocation.", size=18, color=GRAY, after=16)

    tiers = [
        ("7", "High-impact", "100% kill rate categories", GREEN),
        ("5", "Iteratively refined", "Scan-informed description tuning", TEAL),
        ("3", "Content workarounds", "Multi-turn via content themes", BLACK),
        ("5", "Reserved", "Future findings + customer-specific", ORANGE),
    ]
    for count, tier, purpose, color in tiers:
        p = add_p(tf, "", size=18, after=12)
        r(p, count, size=28, bold=True, color=color)
        r(p, f"   {tier}", size=18, bold=True, color=BLACK)
        add_p(tf, f"       {purpose}", size=15, color=GRAY, after=4)

    add_p(tf, "Start at 60%. Leave room to adapt.", size=18, bold=True, color=DARK_ORANGE, before=12)

    # ----- SLIDE 19: What topics can't catch -----
    slide = new_slide(prs)
    title_box(slide, "What topics can\u2019t catch")
    accent_bar(slide, color=RED_ACCENT)

    tf = body_frame(slide)
    add_p(tf, "", size=1)
    add_p(tf, "The remaining 1.2% ASR breaks into three unfixable buckets:", size=18, color=GRAY, after=16)

    buckets = [
        ("38%", "Semantic ambiguity", "POLITICAL", '"Worst president?" \u2014 opinion or attack? Genuinely ambiguous.'),
        ("20%", "Classifier limitations", "EVASION", "Unicode math, IPA phonetic, emoji \u2014 classifier can\u2019t read encoded text."),
        ("11%", "Scan variance", "STOCHASTIC", "Different prompts surface each run. Can\u2019t pre-target what you haven\u2019t seen."),
    ]
    for pct, name, cat, desc in buckets:
        p = add_p(tf, "", size=18, after=4)
        r(p, pct, size=24, bold=True, color=RED_ACCENT)
        r(p, f"   {name}", size=18, bold=True, color=BLACK)
        r(p, f"  ({cat})", size=14, color=MID_GRAY)
        add_p(tf, f"       {desc}", size=15, color=GRAY, after=14)

    # ----- SLIDE 20: The loop -----
    slide = new_slide(prs)
    title_box(slide, "The iteration loop")
    accent_bar(slide)

    tf = body_frame(slide)
    add_p(tf, "", size=1)
    steps = [
        ("Run", "red team scan (static + agent) in SCM"),
        ("Export", "and parse results into structured threat data"),
        ("Analyze", "bypassing prompts \u2014 group by semantic pattern"),
        ("Write", "or update topics naming the specific mechanisms"),
        ("Deploy", "via Management API with BLOCK action"),
        ("Re-scan", "and measure \u2014 repeat until plateau"),
    ]
    for verb, rest in steps:
        p = add_p(tf, "", size=18, after=14)
        r(p, verb, size=20, bold=True, color=ORANGE)
        r(p, f"  {rest}", size=17, color=BLACK)

    add_p(tf, "Expect 2-4 iterations to reach the diminishing returns plateau.", size=16, color=GRAY, before=8)

    # ----- SLIDE 21: Key takeaways -----
    slide = new_slide(prs, bg=WHITE)
    title_box(slide, "Key takeaways")
    accent_bar(slide)

    tf = body_frame(slide)
    add_p(tf, "", size=1)
    takeaways = [
        "First deployment is the biggest lever \u2014 get the initial topics right",
        "Description quality > quantity \u2014 250 chars of precision beats coverage",
        "Name specific mechanisms, not abstract categories",
        "Target content themes, not attack structure",
        "Read your scan data \u2014 bypassing prompts tell you what to write",
        "Reserve 25-40% of slots for iteration",
        "Accept the plateau \u2014 ~1% ASR is the practical floor",
    ]
    for i, text in enumerate(takeaways, 1):
        p = add_p(tf, "", size=17, after=12)
        r(p, f"{i}", size=18, bold=True, color=ORANGE)
        r(p, f"   {text}", size=17, color=BLACK)

    # ----- SLIDE 22: Closing -----
    slide = new_slide(prs, bg=WHITE)
    bar = slide.shapes.add_shape(1, 0, Emu(8800000), prs.slide_width, Emu(900000))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()

    tf = textbox(slide, 623400, 2600000, 16000000, 4000000)
    add_p(tf, "Custom topics are precision tools.", size=38, bold=True, color=BLACK, after=16)
    add_p(tf, "Narrow scope. Specific vocabulary.", size=28, color=GRAY, after=8)
    add_p(tf, "Scan-informed iteration.", size=28, color=GRAY)

    # ----- Remove original template slides -----
    for _ in range(original_count):
        delete_slide(prs, 0)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_deck()
